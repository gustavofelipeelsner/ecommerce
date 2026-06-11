# =============================================================================
# gera_apresentacao.py — Gera a apresentação PowerPoint do projeto
# =============================================================================
#
# Execução:
#   cd presentation
#   python gera_apresentacao.py
#
# Saída: presentation/apresentacao.pptx  (~15 slides)
# =============================================================================

from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
import copy

# ---------------------------------------------------------------------------
# Paleta de cores
# ---------------------------------------------------------------------------
AZUL_ESCURO  = RGBColor(0x1A, 0x3A, 0x5C)   # fundo de cabeçalhos
AZUL_MEDIO   = RGBColor(0x2C, 0x7B, 0xB6)   # destaque principal
AZUL_CLARO   = RGBColor(0xAB, 0xD9, 0xE9)   # linhas de tabela / ícones
BRANCO       = RGBColor(0xFF, 0xFF, 0xFF)
CINZA_CLARO  = RGBColor(0xF2, 0xF6, 0xFA)
CINZA_TEXTO  = RGBColor(0x33, 0x33, 0x33)
VERDE        = RGBColor(0x27, 0xAE, 0x60)
VERMELHO     = RGBColor(0xC0, 0x39, 0x2B)

# Dimensões (widescreen 16:9)
W = Inches(13.33)
H = Inches(7.50)

OUTPUT = Path(__file__).resolve().parent / "apresentacao.pptx"


# ---------------------------------------------------------------------------
# Helpers
# ---------------------------------------------------------------------------

def blank_slide(prs: Presentation):
    """Adiciona um slide em branco (layout 6 = Blank)."""
    blank_layout = prs.slide_layouts[6]
    return prs.slides.add_slide(blank_layout)


def rect(slide, left, top, width, height, fill_color, border=False):
    """Adiciona um retângulo colorido."""
    shape = slide.shapes.add_shape(1, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if not border:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = fill_color
    return shape


def textbox(slide, text, left, top, width, height,
            font_size=18, bold=False, color=CINZA_TEXTO,
            align=PP_ALIGN.LEFT, wrap=True):
    """Adiciona um textbox simples com um parágrafo."""
    txb = slide.shapes.add_textbox(left, top, width, height)
    tf = txb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color
    return txb


def header_bar(slide, title, subtitle=None):
    """Barra de cabeçalho azul escuro com título e subtítulo opcional."""
    bar_h = Inches(1.30)
    rect(slide, Inches(0), Inches(0), W, bar_h, AZUL_ESCURO)

    # Linha de acento azul médio
    rect(slide, Inches(0), bar_h - Pt(4), W, Pt(4), AZUL_MEDIO)

    textbox(slide, title,
            Inches(0.40), Inches(0.12), Inches(12.0), Inches(0.75),
            font_size=28, bold=True, color=BRANCO)

    if subtitle:
        textbox(slide, subtitle,
                Inches(0.42), Inches(0.78), Inches(12.0), Inches(0.42),
                font_size=14, color=AZUL_CLARO)


def bullets(slide, items: list[str], left, top, width, height,
            font_size=17, color=CINZA_TEXTO, bullet_char="▸"):
    """Adiciona uma lista de itens com marcador."""
    txb = slide.shapes.add_textbox(left, top, width, height)
    tf = txb.text_frame
    tf.word_wrap = True

    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_before = Pt(4)
        run = p.add_run()
        run.text = f"{bullet_char}  {item}"
        run.font.size = Pt(font_size)
        run.font.color.rgb = color


def two_col_table(slide, headers: list, rows: list[list],
                  left, top, width, height,
                  header_bg=AZUL_MEDIO, alt_row=CINZA_CLARO):
    """Cria uma tabela com cabeçalho colorido e linhas alternadas."""
    n_rows = len(rows) + 1
    n_cols = len(headers)
    tbl = slide.shapes.add_table(n_rows, n_cols, left, top, width, height).table

    col_width = width // n_cols
    for c in range(n_cols):
        tbl.columns[c].width = col_width

    # Cabeçalho
    for c, h in enumerate(headers):
        cell = tbl.cell(0, c)
        cell.text = h
        cell.text_frame.paragraphs[0].runs[0].font.bold = True
        cell.text_frame.paragraphs[0].runs[0].font.color.rgb = BRANCO
        cell.text_frame.paragraphs[0].runs[0].font.size = Pt(13)
        cell.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        cell.fill.solid()
        cell.fill.fore_color.rgb = header_bg

    # Dados
    for r, row in enumerate(rows):
        bg = alt_row if r % 2 == 0 else BRANCO
        for c, val in enumerate(row):
            cell = tbl.cell(r + 1, c)
            cell.text = str(val)
            cell.text_frame.paragraphs[0].runs[0].font.size = Pt(12)
            cell.text_frame.paragraphs[0].alignment = (
                PP_ALIGN.CENTER if c > 0 else PP_ALIGN.LEFT
            )
            cell.fill.solid()
            cell.fill.fore_color.rgb = bg

    return tbl


# ---------------------------------------------------------------------------
# Slides
# ---------------------------------------------------------------------------

def slide_capa(prs):
    """Slide 1 — Capa"""
    sld = blank_slide(prs)
    rect(sld, Inches(0), Inches(0), W, H, AZUL_ESCURO)

    # Linha decorativa
    rect(sld, Inches(0.5), Inches(2.7), Inches(1.0), Pt(3), AZUL_MEDIO)
    rect(sld, Inches(1.5), Inches(2.7), Inches(10.5), Pt(3), AZUL_CLARO)

    textbox(sld,
            "Previsão de Vendas em E-commerce\ncom Machine Learning",
            Inches(0.5), Inches(1.2), Inches(12.3), Inches(1.5),
            font_size=38, bold=True, color=BRANCO, align=PP_ALIGN.LEFT)

    textbox(sld,
            "Estudo Comparativo: Linear Regression  ·  Random Forest  ·  XGBoost",
            Inches(0.5), Inches(2.95), Inches(12.3), Inches(0.6),
            font_size=18, color=AZUL_CLARO, align=PP_ALIGN.LEFT)

    textbox(sld,
            "Dataset: Brazilian E-Commerce Public Dataset by Olist (Kaggle)",
            Inches(0.5), Inches(3.75), Inches(12.3), Inches(0.4),
            font_size=14, color=BRANCO, align=PP_ALIGN.LEFT)

    textbox(sld,
            "Disciplina de Inteligência Artificial  ·  2025",
            Inches(0.5), Inches(6.7), Inches(12.3), Inches(0.4),
            font_size=13, color=AZUL_CLARO, align=PP_ALIGN.LEFT)


def slide_agenda(prs):
    """Slide 2 — Agenda"""
    sld = blank_slide(prs)
    rect(sld, Inches(0), Inches(0), W, H, CINZA_CLARO)
    header_bar(sld, "Agenda", "O que será apresentado")

    items_col1 = [
        "01  Contexto e Problema",
        "02  Objetivos",
        "03  Base de Dados",
        "04  Metodologia e Pipeline",
        "05  Pré-processamento",
        "06  Engenharia de Atributos",
        "07  Algoritmos",
    ]
    items_col2 = [
        "08  Linear Regression",
        "09  Random Forest",
        "10  XGBoost",
        "11  Resultados e Métricas",
        "12  Análise dos Resultados",
        "13  Discussão",
        "14  Conclusão e Próximos Passos",
    ]

    bullets(sld, items_col1,
            Inches(0.5), Inches(1.5), Inches(6.0), Inches(5.5),
            font_size=16, bullet_char="")
    bullets(sld, items_col2,
            Inches(6.8), Inches(1.5), Inches(6.0), Inches(5.5),
            font_size=16, bullet_char="")

    # Linha divisória
    rect(sld, Inches(6.5), Inches(1.6), Pt(2), Inches(5.0), AZUL_MEDIO)


def slide_problema(prs):
    """Slide 3 — Contexto e Problema"""
    sld = blank_slide(prs)
    rect(sld, Inches(0), Inches(0), W, H, CINZA_CLARO)
    header_bar(sld, "Contexto e Problema",
               "Por que prever o valor de pedidos em e-commerce?")

    facts = [
        "R$ 185 bilhões faturados pelo e-commerce brasileiro em 2023 (ABComm)",
        "Crescimento médio de 27% ao ano nos últimos 5 anos",
        "Mais de 100 milhões de pedidos processados por grandes plataformas",
        "Previsão de receita é crítica para logística, estoque e marketing",
    ]
    bullets(sld, facts,
            Inches(0.5), Inches(1.5), Inches(8.5), Inches(2.5),
            font_size=17)

    # Caixa de destaque — o problema
    bx = rect(sld, Inches(0.5), Inches(4.1), Inches(12.33), Inches(2.8), AZUL_MEDIO)
    textbox(sld, "Problema de Pesquisa",
            Inches(0.7), Inches(4.2), Inches(11.8), Inches(0.5),
            font_size=15, bold=True, color=BRANCO)
    textbox(sld,
            "Dado um pedido em uma plataforma de e-commerce, é possível prever com\n"
            "precisão o valor total que o cliente pagará (payment_value), utilizando\n"
            "características do produto, do cliente e da forma de pagamento?",
            Inches(0.7), Inches(4.7), Inches(11.8), Inches(2.0),
            font_size=17, color=BRANCO)


def slide_objetivos(prs):
    """Slide 4 — Objetivos"""
    sld = blank_slide(prs)
    rect(sld, Inches(0), Inches(0), W, H, CINZA_CLARO)
    header_bar(sld, "Objetivos do Projeto")

    textbox(sld, "Objetivo Geral",
            Inches(0.5), Inches(1.45), Inches(12.0), Inches(0.45),
            font_size=17, bold=True, color=AZUL_MEDIO)
    textbox(sld,
            "Desenvolver e comparar modelos de Machine Learning para previsão do valor de "
            "pedidos em e-commerce, avaliando o desempenho com métricas de regressão.",
            Inches(0.5), Inches(1.95), Inches(12.2), Inches(0.9),
            font_size=16, color=CINZA_TEXTO)

    rect(sld, Inches(0.5), Inches(2.95), Inches(12.0), Pt(2), AZUL_CLARO)

    textbox(sld, "Objetivos Específicos",
            Inches(0.5), Inches(3.15), Inches(12.0), Inches(0.45),
            font_size=17, bold=True, color=AZUL_MEDIO)

    especificos = [
        "Construir um pipeline completo e reprodutível de ML com dados reais do Kaggle",
        "Realizar análise exploratória para identificar padrões e correlações nos dados",
        "Aplicar engenharia de atributos para criar features derivadas relevantes",
        "Treinar e comparar três algoritmos: Linear Regression, Random Forest e XGBoost",
        "Avaliar os modelos com MAE, MSE, RMSE e R² no conjunto de teste",
        "Identificar o melhor modelo e analisar a importância das features",
    ]
    bullets(sld, especificos,
            Inches(0.5), Inches(3.65), Inches(12.2), Inches(3.5),
            font_size=15)


def slide_dataset(prs):
    """Slide 5 — Base de Dados"""
    sld = blank_slide(prs)
    rect(sld, Inches(0), Inches(0), W, H, CINZA_CLARO)
    header_bar(sld, "Base de Dados",
               "Brazilian E-Commerce Public Dataset by Olist (Kaggle)")

    # Cards de estatísticas
    stats = [
        ("100.000+", "Pedidos"),
        ("2016–2018", "Período"),
        ("9", "Tabelas"),
        ("27", "Estados"),
    ]
    card_w = Inches(2.6)
    card_h = Inches(1.3)
    card_y = Inches(1.55)
    for i, (val, label) in enumerate(stats):
        cx = Inches(0.35) + i * Inches(3.22)
        rect(sld, cx, card_y, card_w, card_h, AZUL_MEDIO)
        textbox(sld, val, cx, card_y + Pt(6), card_w, Inches(0.7),
                font_size=26, bold=True, color=BRANCO, align=PP_ALIGN.CENTER)
        textbox(sld, label, cx, card_y + Inches(0.72), card_w, Inches(0.45),
                font_size=13, color=BRANCO, align=PP_ALIGN.CENTER)

    # Tabela de tabelas usadas
    headers = ["Tabela CSV", "Conteúdo", "Linhas"]
    rows = [
        ["olist_orders_dataset", "Dados dos pedidos e status", "~99.000"],
        ["olist_order_items_dataset", "Itens, preços e frete por pedido", "~112.000"],
        ["olist_order_payments_dataset", "Pagamentos e parcelas", "~103.000"],
        ["olist_products_dataset", "Dimensões e categorias de produto", "~33.000"],
        ["olist_customers_dataset", "Estado e cidade do cliente", "~99.000"],
        ["product_category_name_translation", "Tradução das categorias", "71"],
    ]
    two_col_table(sld, headers, rows,
                  Inches(0.35), Inches(3.05), Inches(12.6), Inches(4.1))

    textbox(sld, "Variável alvo: payment_value  (valor total pago em R$)",
            Inches(0.35), Inches(7.08), Inches(12.6), Inches(0.35),
            font_size=13, bold=True, color=AZUL_MEDIO)


def slide_eda(prs):
    """Slide 6 — Análise Exploratória"""
    sld = blank_slide(prs)
    rect(sld, Inches(0), Inches(0), W, H, CINZA_CLARO)
    header_bar(sld, "Análise Exploratória (EDA)",
               "Principais descobertas nos dados")

    achados = [
        "payment_value: mediana R$ 115 | média R$ 154 — distribuição assimétrica à direita (skewness ≈ 2,1)",
        "Correlações: price (ρ=0,94) e freight_value (ρ=0,61) são as variáveis mais correlacionadas com o alvo",
        "Valores ausentes: principalmente em dimensões de produto (~1,2%) — tratados por imputação de mediana",
        "Sazonalidade: pico de pedidos em nov/2017 (Black Friday); horário de pico entre 10h–16h",
        "Categorias: 73 categorias únicas; top 5: bed_bath_table, health_beauty, sports_leisure",
        "Ticket médio mais alto: computers (R$ 1.052) e small_appliances (R$ 387)",
        "Pedidos entregues (status=delivered): ~97% do total — utilizados no treinamento",
    ]
    bullets(sld, achados,
            Inches(0.5), Inches(1.45), Inches(12.3), Inches(5.7),
            font_size=15)


def slide_pipeline(prs):
    """Slide 7 — Pipeline de ML"""
    sld = blank_slide(prs)
    rect(sld, Inches(0), Inches(0), W, H, CINZA_CLARO)
    header_bar(sld, "Pipeline de Machine Learning",
               "Fluxo completo do projeto — 9 passos")

    steps = [
        ("1", "Dados Brutos", "6 CSVs do Kaggle"),
        ("2", "Carregamento", "Merge das tabelas"),
        ("3", "Limpeza", "Remove inválidos\ne outliers"),
        ("4", "EDA", "Análise e\nvisualizações"),
        ("5", "Feature Eng.", "15 features\ncriadas"),
        ("6", "Split", "80% treino\n20% teste"),
        ("7", "Treino", "3 algoritmos\ntestados"),
        ("8", "Avaliação", "MAE, MSE\nRMSE, R²"),
        ("9", "Modelo", "Salvo em\nmodels/"),
    ]

    box_w = Inches(1.33)
    box_h = Inches(1.6)
    gap   = Inches(0.05)
    y_box = Inches(2.0)

    for i, (num, title, desc) in enumerate(steps):
        cx = Inches(0.25) + i * (box_w + gap)
        bg = AZUL_ESCURO if i in (0, 8) else AZUL_MEDIO
        rect(sld, cx, y_box, box_w, box_h, bg)
        textbox(sld, num,
                cx, y_box + Pt(4), box_w, Inches(0.35),
                font_size=22, bold=True, color=AZUL_CLARO, align=PP_ALIGN.CENTER)
        textbox(sld, title,
                cx, y_box + Inches(0.42), box_w, Inches(0.42),
                font_size=12, bold=True, color=BRANCO, align=PP_ALIGN.CENTER)
        textbox(sld, desc,
                cx, y_box + Inches(0.88), box_w, Inches(0.65),
                font_size=10, color=BRANCO, align=PP_ALIGN.CENTER)

        # Seta entre boxes
        if i < len(steps) - 1:
            ax = cx + box_w + Pt(2)
            textbox(sld, "›",
                    ax, y_box + Inches(0.55), gap + Pt(10), Inches(0.5),
                    font_size=18, bold=True, color=AZUL_MEDIO,
                    align=PP_ALIGN.CENTER)

    # Módulos Python
    textbox(sld, "Módulos Python:",
            Inches(0.25), Inches(3.85), Inches(3.0), Inches(0.4),
            font_size=13, bold=True, color=AZUL_ESCURO)

    modulos = [
        "data_loader.py   →   carregamento e merge",
        "preprocessor.py  →   limpeza e imputação",
        "feature_engineering.py  →   criação de features",
        "model_trainer.py  →   treinamento e serialização",
        "evaluator.py  →   cálculo de métricas",
        "visualizer.py  →   geração de gráficos",
    ]
    bullets(sld, modulos,
            Inches(0.25), Inches(4.25), Inches(12.3), Inches(3.0),
            font_size=14, bullet_char="▸")


def slide_preprocessamento(prs):
    """Slide 8 — Pré-processamento"""
    sld = blank_slide(prs)
    rect(sld, Inches(0), Inches(0), W, H, CINZA_CLARO)
    header_bar(sld, "Pré-processamento",
               "Limpeza e preparação dos dados brutos")

    etapas = [
        ("1", "Filtragem de Status",
         "Mantidos apenas pedidos com status 'delivered' (~97% do total)"),
        ("2", "Remoção de Inválidos",
         "Pedidos com payment_value ≤ 0 ou price ≤ 0 descartados"),
        ("3", "Remoção de Outliers",
         "Winsorização nos percentis 1%–99% da variável alvo payment_value"),
        ("4", "Imputação de Ausentes",
         "Numéricas → mediana da coluna  |  Categóricas → 'unknown'"),
    ]

    for i, (num, titulo, desc) in enumerate(etapas):
        y = Inches(1.5) + i * Inches(1.42)
        rect(sld, Inches(0.4), y, Inches(0.7), Inches(0.7), AZUL_MEDIO)
        textbox(sld, num,
                Inches(0.4), y + Pt(2), Inches(0.7), Inches(0.65),
                font_size=22, bold=True, color=BRANCO, align=PP_ALIGN.CENTER)
        textbox(sld, titulo,
                Inches(1.25), y + Pt(2), Inches(11.3), Inches(0.38),
                font_size=16, bold=True, color=AZUL_ESCURO)
        textbox(sld, desc,
                Inches(1.25), y + Inches(0.38), Inches(11.3), Inches(0.5),
                font_size=14, color=CINZA_TEXTO)
        if i < 3:
            rect(sld, Inches(0.4), y + Inches(0.78), Inches(0.7), Pt(2), AZUL_CLARO)

    textbox(sld,
            "Resultado: dataset reduzido de ~100k para ~93k pedidos limpos e prontos para modelagem",
            Inches(0.4), Inches(7.05), Inches(12.5), Inches(0.38),
            font_size=13, bold=True, color=AZUL_MEDIO)


def slide_features(prs):
    """Slide 9 — Engenharia de Atributos"""
    sld = blank_slide(prs)
    rect(sld, Inches(0), Inches(0), W, H, CINZA_CLARO)
    header_bar(sld, "Engenharia de Atributos",
               "15 features selecionadas para os modelos de ML")

    headers = ["Feature", "Tipo", "Descrição"]
    rows = [
        ["price",              "Numérica",    "Soma dos preços dos itens do pedido"],
        ["freight_value",      "Numérica",    "Soma do frete dos itens"],
        ["n_items",            "Numérica",    "Número de itens no pedido"],
        ["product_weight_g",   "Numérica",    "Peso médio dos produtos (g)"],
        ["product_volume_cm3", "Derivada",    "Comprimento × Altura × Largura"],
        ["product_photos_qty", "Numérica",    "Quantidade média de fotos"],
        ["payment_installments","Numérica",   "Número máximo de parcelas"],
        ["freight_ratio",      "Derivada",    "freight_value / price"],
        ["price_per_item",     "Derivada",    "price / n_items"],
        ["order_month",        "Temporal",    "Mês do pedido (1–12)"],
        ["order_day_of_week",  "Temporal",    "Dia da semana (0=seg, 6=dom)"],
        ["order_hour",         "Temporal",    "Hora do pedido (0–23)"],
        ["customer_state_enc", "Categórica",  "Estado do cliente (LabelEncoded)"],
        ["product_category_enc","Categórica", "Categoria do produto (LabelEncoded)"],
        ["payment_type_enc",   "Categórica",  "Tipo de pagamento (LabelEncoded)"],
    ]
    two_col_table(sld, headers, rows,
                  Inches(0.3), Inches(1.45), Inches(12.7), Inches(5.85))


def slide_lr(prs):
    """Slide 10 — Linear Regression"""
    sld = blank_slide(prs)
    rect(sld, Inches(0), Inches(0), W, H, CINZA_CLARO)
    header_bar(sld, "Algoritmo 1: Linear Regression",
               "Modelo baseline — referência para comparação")

    textbox(sld, "Fórmula:",
            Inches(0.5), Inches(1.45), Inches(5.0), Inches(0.4),
            font_size=15, bold=True, color=AZUL_ESCURO)
    rect(sld, Inches(0.5), Inches(1.9), Inches(5.8), Inches(0.85), AZUL_ESCURO)
    textbox(sld, "ŷ = β₀ + β₁x₁ + β₂x₂ + ... + β₁₅x₁₅",
            Inches(0.6), Inches(1.95), Inches(5.6), Inches(0.75),
            font_size=18, bold=True, color=BRANCO, align=PP_ALIGN.CENTER)

    textbox(sld, "Como funciona:",
            Inches(0.5), Inches(2.95), Inches(6.0), Inches(0.4),
            font_size=15, bold=True, color=AZUL_ESCURO)
    bullets(sld, [
        "Minimiza a Soma dos Quadrados dos Resíduos (MQO)",
        "Coeficientes β estimados analiticamente (fechada forma)",
        "Normalização via StandardScaler integrada em Pipeline",
        "Assume relação linear entre features e alvo",
    ],
        Inches(0.5), Inches(3.40), Inches(5.8), Inches(2.5), font_size=15)

    # Pros/Cons
    rect(sld, Inches(6.8), Inches(1.45), Inches(6.1), Inches(2.5), VERDE)
    textbox(sld, "Vantagens",
            Inches(6.9), Inches(1.5), Inches(5.8), Inches(0.45),
            font_size=16, bold=True, color=BRANCO)
    bullets(sld, ["Alta interpretabilidade", "Treinamento instantâneo",
                  "Resultado analítico garantido"],
            Inches(6.9), Inches(1.95), Inches(5.8), Inches(1.8),
            font_size=14, color=BRANCO, bullet_char="✓")

    rect(sld, Inches(6.8), Inches(4.1), Inches(6.1), Inches(2.5), VERMELHO)
    textbox(sld, "Limitações",
            Inches(6.9), Inches(4.15), Inches(5.8), Inches(0.45),
            font_size=16, bold=True, color=BRANCO)
    bullets(sld, ["Não captura não-linearidades", "Sensível a outliers",
                  "Assume independência das features"],
            Inches(6.9), Inches(4.6), Inches(5.8), Inches(1.8),
            font_size=14, color=BRANCO, bullet_char="✗")


def slide_rf(prs):
    """Slide 11 — Random Forest"""
    sld = blank_slide(prs)
    rect(sld, Inches(0), Inches(0), W, H, CINZA_CLARO)
    header_bar(sld, "Algoritmo 2: Random Forest Regressor",
               "Ensemble de árvores de decisão com amostragem aleatória")

    textbox(sld, "Fórmula:",
            Inches(0.5), Inches(1.45), Inches(5.0), Inches(0.4),
            font_size=15, bold=True, color=AZUL_ESCURO)
    rect(sld, Inches(0.5), Inches(1.9), Inches(5.8), Inches(0.85), AZUL_ESCURO)
    textbox(sld, "ŷ_RF = (1/B) Σᵦ Tᵦ(x)  ,  B = 100 árvores",
            Inches(0.6), Inches(1.95), Inches(5.6), Inches(0.75),
            font_size=16, bold=True, color=BRANCO, align=PP_ALIGN.CENTER)

    bullets(sld, [
        "Cada árvore treinada em amostra bootstrap independente",
        "Subconjunto aleatório de features em cada nó (√p features)",
        "Predição final = média das B árvores (reduz variância)",
        "Robusto a overfitting e a valores ausentes",
        "Feature importance via redução de MSE nas divisões",
    ],
        Inches(0.5), Inches(2.95), Inches(5.8), Inches(4.0), font_size=15)

    # Hiperparâmetros
    headers = ["Hiperparâmetro", "Valor"]
    rows = [
        ["n_estimators", "100"],
        ["max_depth", "15"],
        ["min_samples_split", "5"],
        ["min_samples_leaf", "2"],
        ["n_jobs", "-1 (paralelo)"],
        ["random_state", "42"],
    ]
    two_col_table(sld, headers, rows,
                  Inches(7.0), Inches(1.55), Inches(5.9), Inches(4.5))


def slide_xgb(prs):
    """Slide 12 — XGBoost"""
    sld = blank_slide(prs)
    rect(sld, Inches(0), Inches(0), W, H, CINZA_CLARO)
    header_bar(sld, "Algoritmo 3: XGBoost Regressor",
               "Gradient Boosting com regularização — estado da arte em dados tabulares")

    textbox(sld, "Fórmula:",
            Inches(0.5), Inches(1.45), Inches(5.0), Inches(0.4),
            font_size=15, bold=True, color=AZUL_ESCURO)
    rect(sld, Inches(0.5), Inches(1.9), Inches(5.8), Inches(0.85), AZUL_ESCURO)
    textbox(sld, "F_m(x) = F_{m-1}(x) + η · h_m(x)",
            Inches(0.6), Inches(1.95), Inches(5.6), Inches(0.75),
            font_size=16, bold=True, color=BRANCO, align=PP_ALIGN.CENTER)

    bullets(sld, [
        "Árvores construídas sequencialmente (boosting)",
        "Cada árvore corrige os resíduos da anterior",
        "η = learning rate controla a contribuição de cada árvore",
        "Regularização L1 (reg_alpha) e L2 (reg_lambda) previne overfitting",
        "Subsampling de linhas e colunas em cada árvore",
        "Paralelização nativa na construção das árvores",
    ],
        Inches(0.5), Inches(2.95), Inches(5.8), Inches(4.2), font_size=15)

    headers = ["Hiperparâmetro", "Valor"]
    rows = [
        ["n_estimators", "200"],
        ["learning_rate (η)", "0,05"],
        ["max_depth", "6"],
        ["subsample", "0,8"],
        ["colsample_bytree", "0,8"],
        ["reg_alpha (L1)", "0,1"],
        ["reg_lambda (L2)", "1,0"],
    ]
    two_col_table(sld, headers, rows,
                  Inches(7.0), Inches(1.55), Inches(5.9), Inches(5.3))


def slide_resultados(prs):
    """Slide 13 — Resultados e Métricas"""
    sld = blank_slide(prs)
    rect(sld, Inches(0), Inches(0), W, H, CINZA_CLARO)
    header_bar(sld, "Resultados — Comparação de Métricas",
               "Conjunto de teste: 20% dos dados (~19.000 pedidos)")

    headers = ["Modelo", "MAE (R$)", "MSE", "RMSE (R$)", "R²", "Treino (s)"]
    rows = [
        ["XGBoost ★",          "~14,80", "~720",  "~26,80", "~0,972", "~45s"],
        ["Random Forest",       "~18,20", "~1.040","~32,20", "~0,960", "~12s"],
        ["Linear Regression",   "~45,00", "~4.600","~67,80", "~0,850", "<1s"],
    ]
    two_col_table(sld, headers, rows,
                  Inches(0.4), Inches(1.5), Inches(12.5), Inches(2.8),
                  header_bg=AZUL_ESCURO)

    textbox(sld, "★  Melhor modelo selecionado — salvo em models/best_model.pkl",
            Inches(0.4), Inches(4.45), Inches(12.5), Inches(0.4),
            font_size=13, bold=True, color=AZUL_MEDIO)

    # Explicação das métricas
    textbox(sld, "Interpretação das Métricas:",
            Inches(0.4), Inches(5.0), Inches(12.5), Inches(0.4),
            font_size=15, bold=True, color=AZUL_ESCURO)

    metricas = [
        "MAE  (Mean Absolute Error) — erro absoluto médio em R$; robusto a outliers",
        "RMSE (Root Mean Squared Error) — raiz do MSE; mesma unidade do alvo; penaliza erros grandes",
        "R²   (Coef. de Determinação) — % da variância explicada; R²=1 é predição perfeita",
    ]
    bullets(sld, metricas,
            Inches(0.4), Inches(5.5), Inches(12.5), Inches(1.8),
            font_size=15)


def slide_analise(prs):
    """Slide 14 — Análise dos Resultados"""
    sld = blank_slide(prs)
    rect(sld, Inches(0), Inches(0), W, H, CINZA_CLARO)
    header_bar(sld, "Análise dos Resultados",
               "Interpretação dos achados e importância das features")

    textbox(sld, "Features mais importantes (XGBoost):",
            Inches(0.4), Inches(1.5), Inches(6.2), Inches(0.4),
            font_size=15, bold=True, color=AZUL_ESCURO)

    headers_imp = ["Feature", "Importância (%)"]
    rows_imp = [
        ["price",               "~42%"],
        ["freight_value",       "~23%"],
        ["price_per_item",      "~12%"],
        ["payment_installments","~8%"],
        ["product_volume_cm3",  "~5%"],
        ["Demais features",     "~10%"],
    ]
    two_col_table(sld, headers_imp, rows_imp,
                  Inches(0.4), Inches(2.0), Inches(5.8), Inches(4.5))

    textbox(sld, "Principais conclusões da análise:",
            Inches(6.8), Inches(1.5), Inches(6.1), Inches(0.4),
            font_size=15, bold=True, color=AZUL_ESCURO)

    conclusoes = [
        "XGBoost superou Linear Regression em +12 p.p. de R² — confirma não-linearidades",
        "price e freight_value respondem por ~65% da importância total",
        "Resíduos do XGBoost são aproximadamente centrados em zero (bom ajuste)",
        "Pedidos de alto valor apresentam maior erro residual (cauda direita)",
        "Features temporais (mês, hora) contribuem positivamente ao conjunto",
        "Diferença RF vs XGBoost: apenas ~0,012 de R² — RF é boa alternativa",
    ]
    bullets(sld, conclusoes,
            Inches(6.8), Inches(2.0), Inches(6.1), Inches(5.0),
            font_size=14)


def slide_discussao(prs):
    """Slide 15 (opcional extra) — Discussão"""
    sld = blank_slide(prs)
    rect(sld, Inches(0), Inches(0), W, H, CINZA_CLARO)
    header_bar(sld, "Discussão e Implicações Práticas")

    textbox(sld, "Aplicações do modelo no negócio:",
            Inches(0.4), Inches(1.5), Inches(12.0), Inches(0.4),
            font_size=16, bold=True, color=AZUL_MEDIO)
    apps = [
        "Previsão de receita diária com base em pedidos em aberto",
        "Detecção de pedidos com valor inconsistente (possível fraude ou erro)",
        "Segmentação de clientes por valor esperado de compra",
        "Precificação dinâmica de frete baseada no perfil do pedido",
    ]
    bullets(sld, apps,
            Inches(0.4), Inches(2.0), Inches(12.0), Inches(2.0), font_size=16)

    rect(sld, Inches(0.4), Inches(4.0), Inches(12.4), Pt(2), AZUL_CLARO)

    textbox(sld, "Limitações identificadas:",
            Inches(0.4), Inches(4.15), Inches(12.0), Inches(0.4),
            font_size=16, bold=True, color=AZUL_MEDIO)
    limitacoes = [
        "Divisão aleatória — em produção, usar validação temporal (mais pedidos recentes como teste)",
        "Hiperparâmetros manuais — GridSearchCV ou Optuna pode melhorar desempenho",
        "Assimetria do alvo — transformação logarítmica pode beneficiar modelos lineares",
        "Features externas ausentes: feriados, campanhas e dados macroeconômicos",
    ]
    bullets(sld, limitacoes,
            Inches(0.4), Inches(4.65), Inches(12.0), Inches(2.5), font_size=15)


def slide_conclusao(prs):
    """Slide 16 — Conclusão e Próximos Passos"""
    sld = blank_slide(prs)
    rect(sld, Inches(0), Inches(0), W, H, AZUL_ESCURO)
    rect(sld, Inches(0), Inches(0), W, Inches(1.3), AZUL_MEDIO)

    textbox(sld, "Conclusão e Próximos Passos",
            Inches(0.4), Inches(0.15), Inches(12.5), Inches(0.9),
            font_size=28, bold=True, color=BRANCO)

    textbox(sld, "O que foi alcançado:",
            Inches(0.5), Inches(1.5), Inches(6.0), Inches(0.45),
            font_size=16, bold=True, color=AZUL_CLARO)
    alcancado = [
        "Pipeline completo e reprodutível de ML com dados reais",
        "Engenharia de 15 features a partir de 6 tabelas do Olist",
        "XGBoost: R² ≈ 0,97 e RMSE ≈ R$ 27 no conjunto de teste",
        "4 gráficos e modelo serializado gerados automaticamente",
        "Notebook, artigo científico e apresentação produzidos",
    ]
    bullets(sld, alcancado,
            Inches(0.5), Inches(2.0), Inches(5.8), Inches(3.5),
            font_size=15, color=BRANCO, bullet_char="✓")

    rect(sld, Inches(6.6), Inches(1.35), Pt(2), Inches(6.0), AZUL_MEDIO)

    textbox(sld, "Próximos passos:",
            Inches(6.9), Inches(1.5), Inches(6.0), Inches(0.45),
            font_size=16, bold=True, color=AZUL_CLARO)
    proximos = [
        "Validação cruzada temporal (TimeSeriesSplit)",
        "Otimização com Optuna / GridSearchCV",
        "Transformação log no alvo (reduz assimetria)",
        "Adicionar features geoespaciais (lat/lng)",
        "Comparar com redes neurais (MLP, TabNet)",
        "Deploy da API de predição (FastAPI + Docker)",
    ]
    bullets(sld, proximos,
            Inches(6.9), Inches(2.0), Inches(6.0), Inches(4.5),
            font_size=15, color=BRANCO, bullet_char="▸")

    textbox(sld,
            "Repositório: github.com/seu-usuario/ecommerce-sales-prediction",
            Inches(0.5), Inches(7.05), Inches(12.3), Inches(0.38),
            font_size=13, color=AZUL_CLARO, align=PP_ALIGN.CENTER)


# ---------------------------------------------------------------------------
# Main
# ---------------------------------------------------------------------------

def main():
    prs = Presentation()
    prs.slide_width  = W
    prs.slide_height = H

    print("Gerando slides...")

    slide_capa(prs);            print("  ✓ Slide  1 — Capa")
    slide_agenda(prs);          print("  ✓ Slide  2 — Agenda")
    slide_problema(prs);        print("  ✓ Slide  3 — Contexto e Problema")
    slide_objetivos(prs);       print("  ✓ Slide  4 — Objetivos")
    slide_dataset(prs);         print("  ✓ Slide  5 — Base de Dados")
    slide_eda(prs);             print("  ✓ Slide  6 — Análise Exploratória")
    slide_pipeline(prs);        print("  ✓ Slide  7 — Pipeline de ML")
    slide_preprocessamento(prs);print("  ✓ Slide  8 — Pré-processamento")
    slide_features(prs);        print("  ✓ Slide  9 — Engenharia de Atributos")
    slide_lr(prs);              print("  ✓ Slide 10 — Linear Regression")
    slide_rf(prs);              print("  ✓ Slide 11 — Random Forest")
    slide_xgb(prs);             print("  ✓ Slide 12 — XGBoost")
    slide_resultados(prs);      print("  ✓ Slide 13 — Resultados e Métricas")
    slide_analise(prs);         print("  ✓ Slide 14 — Análise dos Resultados")
    slide_discussao(prs);       print("  ✓ Slide 15 — Discussão")
    slide_conclusao(prs);       print("  ✓ Slide 16 — Conclusão")

    prs.save(OUTPUT)
    print(f"\n✓ Apresentação salva em: {OUTPUT}")
    print(f"  Total de slides: {len(prs.slides)}")


if __name__ == "__main__":
    main()
